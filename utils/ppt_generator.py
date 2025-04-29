import os
import requests
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from utils.ai_helpers import generate_summary_and_significance  # (you need to define this function if not done yet)
from utils.confluence import get_image_url
from requests.auth import HTTPBasicAuth
from utils.confluence import get_image_url
from utils.helpers import download_and_encode_image


CONFLUENCE_EMAIL = os.getenv("CONFLUENCE_EMAIL")
CONFLUENCE_API_TOKEN = os.getenv("CONFLUENCE_API_TOKEN")

auth = HTTPBasicAuth(CONFLUENCE_EMAIL, CONFLUENCE_API_TOKEN)

def build_pptx(feature_title, content, page_id, overview_image=None):#change here
    # 1. Load the PPTX Template
    template_path = "monthly_release_template.pptx"  # <-- make sure this file exists
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # 2. Generate AI texts
    summary_text, improvement_text = generate_summary_and_significance(content) #change this

    from pptx.util import Pt

    # Feature Title - Textbox 9
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # Feature name
        if shape.name == "TextBox 9":
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = feature_title
            run.font.name = "Aptos"
            run.font.size = Pt(24)
            run.font.bold = True

        # Summary block
        if shape.name == "TextBox 11":
            shape.text_frame.clear()

            # Heading: Summary
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Summary\n"
            run.font.name = "Aptos"
            run.font.size = Pt(18)
            run.font.bold = True

            # Body
            p = shape.text_frame.add_paragraph()
            run = p.add_run()
            run.text = summary_text
            run.font.name = "Aptos"
            run.font.size = Pt(14)
            run.font.bold = False

        # Improvement Significance block
        if shape.name == "TextBox 12":
            shape.text_frame.clear()

            # Heading
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Improvement Significance\n"
            run.font.name = "Aptos"
            run.font.size = Pt(18)
            run.font.bold = True

            # Body
            p = shape.text_frame.add_paragraph()
            run = p.add_run()
            run.text = improvement_text
            run.font.name = "Aptos"
            run.font.size = Pt(14)
            run.font.bold = False

    # Insert the image if it exists
    if overview_image:
        img_url = get_image_url(page_id, overview_image)
        encoded_img = download_and_encode_image(img_url)

        if encoded_img:
            # Save the image temporarily
            img_path = "temp_slide_image.png"
            with open(img_path, "wb") as f:
                import base64
                f.write(base64.b64decode(encoded_img.split(",")[1]))

            # 📍 Insert image with correct position
            slide = prs.slides[0]  # Assuming first slide
            left = Inches(7.87)
            top = Inches(1.93)
            width = Inches(4.0)  # Max width
            slide.shapes.add_picture(img_path, left, top, width=width)

            # Clean up temp image
            os.remove(img_path)


    # 5. Save the new PPTX
    output_pptx = "monthly_release_update.pptx"
    prs.save(output_pptx)

    return output_pptx

